from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


PROJECT_ROOT = Path(__file__).resolve().parent
OUTPUT = PROJECT_ROOT / "Federated_Healthcare_Analyzer_Presentation.pptx"
SCREENSHOT = Path(
    r"C:\Users\prath\.cursor\projects\c-Users-prath-Desktop-code\assets\c__Users_prath_AppData_Roaming_Cursor_User_workspaceStorage_43e3b676eb9aab1e28a521d759f23fa8_images_image-cbe37ce3-d114-402e-a8ef-307a97da545d.png"
)


def set_title(slide, title_text, subtitle_text=None):
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.3), Inches(12.0), Inches(0.8))
    tf = title_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(34)
    p.font.bold = True
    p.font.color.rgb = RGBColor(26, 32, 44)
    if subtitle_text:
        sub_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.1), Inches(12.0), Inches(0.5))
        sf = sub_box.text_frame
        sf.clear()
        sp = sf.paragraphs[0]
        sp.text = subtitle_text
        sp.font.size = Pt(16)
        sp.font.color.rgb = RGBColor(74, 85, 104)


def add_bullets(slide, x, y, w, h, lines, font_size=22):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    for idx, line in enumerate(lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = line
        p.level = 0
        p.font.size = Pt(font_size)
        p.font.color.rgb = RGBColor(45, 55, 72)


def add_section_divider(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(15, 23, 42)
    bg.line.fill.background()
    t = slide.shapes.add_textbox(Inches(0.8), Inches(2.5), Inches(11.8), Inches(1.2))
    tf = t.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(226, 232, 240)
    p.alignment = PP_ALIGN.CENTER
    s = slide.shapes.add_textbox(Inches(0.8), Inches(3.7), Inches(11.8), Inches(0.8))
    sf = s.text_frame
    sp = sf.paragraphs[0]
    sp.text = subtitle
    sp.font.size = Pt(20)
    sp.font.color.rgb = RGBColor(148, 163, 184)
    sp.alignment = PP_ALIGN.CENTER


def create_presentation():
    prs = Presentation()

    # Slide 1
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_title(
        slide,
        "Federated Healthcare Analyzer",
        "Graduate Project Presentation",
    )
    add_bullets(
        slide,
        0.9,
        2.0,
        11.5,
        3.8,
        [
            "A privacy-aware federated learning workflow for healthcare classification.",
            "Built with Flower, PyTorch, Streamlit, and local LLM-based interpretation.",
            "Focus: end-to-end training loop, monitoring dashboard, and explainability chat.",
        ],
    )

    # Slide 2
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_title(slide, "Motivation & Problem Statement")
    add_bullets(
        slide,
        0.8,
        1.7,
        12.0,
        4.8,
        [
            "Centralized ML often conflicts with healthcare privacy and governance constraints.",
            "Federated learning keeps data local and shares only model parameters/updates.",
            "Goal: create a reproducible FL prototype with transparent metrics and professor-ready demos.",
            "Challenge: combine ML training, observability, and concise interpretation in one tool.",
        ],
    )

    # Slide 3
    add_section_divider(prs, "System Design", "Architecture and execution flow")

    # Slide 4
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_title(slide, "Architecture Overview")
    add_bullets(
        slide,
        0.8,
        1.8,
        12.0,
        4.8,
        [
            "Flower Server (`federated/server.py`) orchestrates rounds and FedAvg aggregation.",
            "Three FL Clients (`federated/client.py`) train locally on data partitions.",
            "PyTorch model (`model/model.py`) performs binary classification.",
            "Dashboard (`dashboard/app.py`) reads `metrics.json` and visualizes progress.",
            "LLM module (`llm/llm_generator.py`) powers insight generation and follow-up chat.",
        ],
    )

    # Slide 5
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_title(slide, "Execution Flow")
    add_bullets(
        slide,
        0.8,
        1.8,
        12.0,
        4.8,
        [
            "1) Start server at `127.0.0.1:8080` with FedAvg strategy.",
            "2) Launch clients (`client.py 0/1/2`), each trains for local epochs.",
            "3) Clients send updates and evaluation metrics back to server.",
            "4) Server computes weighted aggregates and appends to `metrics.json`.",
            "5) Streamlit dashboard renders KPIs/charts and enables LLM Q&A.",
        ],
    )

    # Slide 6
    add_section_divider(prs, "Code Walkthrough", "Core modules and implementation logic")

    # Slide 7
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_title(slide, "Server Logic (`federated/server.py`)")
    add_bullets(
        slide,
        0.8,
        1.8,
        12.0,
        4.8,
        [
            "Defines `weighted_average(metrics)` for round-level aggregation.",
            "Computes weighted accuracy/loss based on client sample counts.",
            "Writes structured metrics: round, accuracy, loss, clients, examples, timestamp.",
            "Persists to project-level `metrics.json` using stable `Path` resolution.",
            "Runs 5 federated rounds via Flower `ServerConfig(num_rounds=5)`.",
        ],
    )

    # Slide 8
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_title(slide, "Client + Data + Model Logic")
    add_bullets(
        slide,
        0.8,
        1.6,
        12.0,
        5.2,
        [
            "Data (`federated/utils.py`): loads sklearn breast-cancer data, stratified split, standardization.",
            "Partitioning: `np.array_split` distributes training data across 3 clients.",
            "Model (`model/model.py`): MLP (input -> 32 -> 16 -> 2) with ReLU activations.",
            "Training: local epochs with Adam optimizer and CrossEntropyLoss.",
            "Evaluation: each client reports test loss + accuracy through `NumPyClient.evaluate`.",
        ],
    )

    # Slide 9
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_title(slide, "Dashboard & Professional UI Enhancements")
    add_bullets(
        slide,
        0.8,
        1.7,
        12.0,
        5.0,
        [
            "Professional Streamlit layout with sidebar controls and styled KPI cards.",
            "Interactive Plotly charts for accuracy/loss trends across rounds.",
            "Round-level data table with timestamps and sample/client counts.",
            "Model assessment badge and convergence cue for quick interpretation.",
            "Robust states: handles missing, empty, and malformed `metrics.json` gracefully.",
        ],
    )

    # Slide 10
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_title(slide, "AI Insight + Follow-up Chat")
    add_bullets(
        slide,
        0.8,
        1.7,
        12.0,
        4.8,
        [
            "Uses local Ollama chat API (`/api/chat`) with default model `qwen3:4b`.",
            "First response summarizes run quality and suggests next optimization steps.",
            "Session-state conversation enables iterative follow-up questions.",
            "Grounding strategy: metrics + runtime dataset context are injected programmatically.",
            "Prompting rules reduce hallucinations and avoid unwanted disclaimer text.",
        ],
    )

    # Slide 11
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_title(slide, "Demo Snapshot")
    if SCREENSHOT.exists():
        slide.shapes.add_picture(str(SCREENSHOT), Inches(0.6), Inches(1.4), width=Inches(12.1))
        cap = slide.shapes.add_textbox(Inches(0.7), Inches(6.8), Inches(12.0), Inches(0.4))
        tf = cap.text_frame
        p = tf.paragraphs[0]
        p.text = "Dashboard chat view: generated insight and follow-up question flow."
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(74, 85, 104)
    else:
        add_bullets(
            slide,
            0.8,
            2.2,
            12.0,
            3.0,
            ["Screenshot not found in workspace at generation time."],
        )

    # Slide 12
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_title(slide, "Results, Limitations, and Future Work")
    add_bullets(
        slide,
        0.8,
        1.7,
        12.0,
        4.8,
        [
            "Observed convergence: accuracy increases across rounds with decreasing loss.",
            "Current setup is a simulation (single machine, fixed client count, shared test set).",
            "Future: non-IID splits, secure aggregation, differential privacy, richer model tuning.",
            "Potential extension: CI/testing, experiment tracking, and deployment packaging.",
        ],
    )

    prs.save(OUTPUT)
    print(f"Created: {OUTPUT}")


if __name__ == "__main__":
    create_presentation()
